from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

import numpy as np
import pandas as pd


@dataclass(frozen=True)
class ChannelResult:
    normalized: np.ndarray | None
    status: str
    minimum: float
    maximum: float
    mean: float
    std: float
    snr_metric: float


def normalize_channel(
    image: np.ndarray,
    channel_name: str,
    *,
    snr_threshold: float = 2.5,
    lower_sigma: float = 1.2,
) -> ChannelResult:
    """Validate and contrast-normalize one 2-D microscopy channel."""
    array = np.asarray(image, dtype=np.float32)
    if array.ndim != 2:
        raise ValueError(
            f"{channel_name} must be a 2-D image after projection; got shape {array.shape}."
        )

    finite = np.isfinite(array)
    if not finite.any():
        return ChannelResult(None, f"{channel_name} contains no finite pixels.", 0, 0, 0, 0, 0)

    values = array[finite]
    image_min = float(values.min())
    image_max = float(values.max())
    image_mean = float(values.mean())
    image_std = float(values.std())

    if image_max <= image_min:
        return ChannelResult(
            None,
            f"Channel {channel_name} has zero variance (blank).",
            image_min,
            image_max,
            image_mean,
            image_std,
            0.0,
        )

    snr_metric = (image_max - image_mean) / image_std if image_std > 0 else 0.0
    if snr_metric < snr_threshold:
        return ChannelResult(
            None,
            f"Channel {channel_name} lacks signal (SNR metric {snr_metric:.2f} < {snr_threshold:.2f}).",
            image_min,
            image_max,
            image_mean,
            image_std,
            snr_metric,
        )

    lower = image_mean + lower_sigma * image_std
    upper = image_max
    if upper <= lower:
        lower = image_min

    denominator = upper - lower
    if denominator <= 0:
        return ChannelResult(
            None,
            f"Channel {channel_name} could not be normalized.",
            image_min,
            image_max,
            image_mean,
            image_std,
            snr_metric,
        )

    clipped = np.clip(array, lower, upper)
    normalized = (clipped - lower) / denominator
    normalized[~finite] = 0.0

    return ChannelResult(
        normalized.astype(np.float32),
        "Success",
        image_min,
        image_max,
        image_mean,
        image_std,
        snr_metric,
    )


def max_project_channels(image: np.ndarray, channel_axis: int | None = None) -> np.ndarray:
    """
    Convert common microscopy arrays to (C, Y, X).

    Supported common layouts include:
    - (C, Y, X)
    - (Z, C, Y, X)
    - (T, Z, C, Y, X)

    For ambiguous arrays, pass channel_axis explicitly.
    """
    array = np.asarray(image)

    if array.ndim == 2:
        return array[np.newaxis, ...]

    if channel_axis is not None:
        axis = channel_axis % array.ndim
        moved = np.moveaxis(array, axis, 0)
        if moved.ndim == 3:
            return moved
        projection_axes = tuple(range(1, moved.ndim - 2))
        return np.max(moved, axis=projection_axes) if projection_axes else moved

    if array.ndim == 3:
        return array

    if array.ndim == 4:
        # Default expected layout: Z, C, Y, X
        return np.max(array, axis=0)

    if array.ndim == 5:
        # Default expected layout: T, Z, C, Y, X
        return np.max(array, axis=(0, 1))

    raise ValueError(
        f"Unsupported image shape {array.shape}. Supply a channel axis for non-standard data."
    )


def overlap_metrics(mask_a: np.ndarray, mask_b: np.ndarray) -> dict[str, float | int]:
    """Return union-based overlap and directional overlap percentages."""
    a = np.asarray(mask_a, dtype=bool)
    b = np.asarray(mask_b, dtype=bool)
    if a.shape != b.shape:
        raise ValueError(f"Mask shapes differ: {a.shape} versus {b.shape}.")

    a_pixels = int(a.sum())
    b_pixels = int(b.sum())
    overlap_pixels = int(np.logical_and(a, b).sum())
    union_pixels = int(np.logical_or(a, b).sum())

    return {
        "Channel_A_Positive_Pixels": a_pixels,
        "Channel_B_Positive_Pixels": b_pixels,
        "Overlap_Pixels": overlap_pixels,
        "Union_Pixels": union_pixels,
        "Overlay_Percentage": (overlap_pixels / union_pixels * 100.0) if union_pixels else 0.0,
        "A_Overlap_Percentage": (overlap_pixels / a_pixels * 100.0) if a_pixels else 0.0,
        "B_Overlap_Percentage": (overlap_pixels / b_pixels * 100.0) if b_pixels else 0.0,
    }


def create_qc_powerpoint(
    plot_paths: Sequence[Path],
    output_file: Path,
    *,
    title_prefix: str = "QC Review Frame",
) -> None:
    """Create a black-background widescreen PowerPoint containing one plot per slide."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError(
            "PowerPoint export requires python-pptx. Install it with: "
            "python -m pip install python-pptx"
        ) from exc

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]

    for plot_path in sorted(plot_paths):
        slide = presentation.slides.add_slide(blank_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        slide.shapes.add_picture(
            str(plot_path),
            Inches(0.166),
            Inches(1.55),
            width=Inches(13.0),
        )

        textbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8)
        )
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.text = f"{title_prefix}: {plot_path.name.replace('layout_', '')}"
        paragraph.font.bold = True
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    output_file.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_file)


def save_reports(records: Iterable[dict], csv_file: Path, xlsx_file: Path | None = None) -> pd.DataFrame:
    """Save the processing log as CSV and optionally Excel."""
    frame = pd.DataFrame(list(records))
    csv_file.parent.mkdir(parents=True, exist_ok=True)
    frame.to_csv(csv_file, index=False)
    if xlsx_file is not None:
        frame.to_excel(xlsx_file, index=False)
    return frame
